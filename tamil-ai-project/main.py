import os
import tempfile
from io import BytesIO
from datetime import datetime

from flask import Flask, request, jsonify, send_file

from faster_whisper import WhisperModel
from docx import Document
import openpyxl
from pptx import Presentation

app = Flask(__name__)

@app.route("/health", methods=["GET"], endpoint="health_check")
def health_check():
    return jsonify({"status": "ok"})


# ----------------------------
# Config (env vars)
# ----------------------------
WHISPER_MODEL = os.getenv("WHISPER_MODEL", "base")     # tiny/base/small/medium
WHISPER_DEVICE = os.getenv("WHISPER_DEVICE", "cpu")    # cpu for Render/Railway
WHISPER_COMPUTE = os.getenv("WHISPER_COMPUTE", "int8") # int8 recommended for CPU
PORT = int(os.getenv("PORT", "8080"))

# Load Whisper once (important for performance)
whisper = WhisperModel(WHISPER_MODEL, device=WHISPER_DEVICE, compute_type=WHISPER_COMPUTE)

# ----------------------------
# Helpers
# ----------------------------
def transcribe_upload(file_storage) -> dict:
    audio_bytes = file_storage.read()
    filename = file_storage.filename or "audio.webm"

    suffix = ".webm"
    if "." in filename:
        suffix = "." + filename.split(".")[-1].lower()

    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(audio_bytes)
        tmp_path = tmp.name

    try:
        segments, info = whisper.transcribe(
            tmp_path,
            language="ta",      # force Tamil (remove for auto-detect)
            vad_filter=True
        )
        text = " ".join(seg.text.strip() for seg in segments).strip()
        return {"text": text, "language": info.language, "duration": getattr(info, "duration", None)}
    finally:
        try:
            os.remove(tmp_path)
        except Exception:
            pass

def detect_intent(text: str) -> dict:
    t = (text or "").strip()
    tl = t.lower()

    if any(w in t for w in ["வேலை", "டாஸ்க்"]) or "task" in tl:
        return {"intent": "ADD_TASK", "params": {"task": t}}

    if any(w in t for w in ["குறிப்பு", "நோட்", "நோட்டு"]) or "note" in tl:
        return {"intent": "ADD_NOTE", "params": {"note": t}}

    if "excel" in tl or "எக்சல்" in t:
        return {"intent": "GENERATE_EXCEL", "params": {"sheet": "Report"}}

    if "ppt" in tl or "powerpoint" in tl or "பவர்" in t:
        return {"intent": "GENERATE_PPT", "params": {"title": "Summary"}}

    if "word" in tl or "doc" in tl or "வேர்ட்" in t:
        return {"intent": "GENERATE_WORD", "params": {"title": "Document"}}

    return {"intent": "UNKNOWN", "params": {"raw": t}}

def make_word_bytes(title: str, content: str) -> bytes:
    doc = Document()
    doc.add_heading(title or "Document", level=1)
    doc.add_paragraph(content or "")
    doc.add_paragraph(f"Generated: {datetime.now().isoformat(sep=' ', timespec='seconds')}")
    bio = BytesIO()
    doc.save(bio)
    return bio.getvalue()

def make_excel_bytes(sheet_name: str, rows: list[list[str]]) -> bytes:
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = (sheet_name or "Sheet1")[:31]
    for r in rows:
        ws.append(r)
    bio = BytesIO()
    wb.save(bio)
    return bio.getvalue()

def make_ppt_bytes(title: str, bullets: list[str]) -> bytes:
    prs = Presentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title or "Presentation"

    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    bullets = bullets or ["(empty)"]
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b

    bio = BytesIO()
    prs.save(bio)
    return bio.getvalue()

# ----------------------------
# Routes
# ----------------------------
@app.get("/")
def home():
    return "Tamil Voice-to-Task AI is Active!"

@app.get("/health")
def health():
    return jsonify({"status": "ok"})

@app.post("/transcribe")
def transcribe():
    if "audio" not in request.files:
        return jsonify({"error": "Missing file field 'audio'"}), 400
    stt = transcribe_upload(request.files["audio"])
    return jsonify(stt)

@app.post("/intent")
def intent():
    data = request.get_json(silent=True) or {}
    text = data.get("text", "")
    return jsonify(detect_intent(text))

@app.post("/pipeline")
def pipeline():
    if "audio" not in request.files:
        return jsonify({"error": "Missing file field 'audio'"}), 400

    stt = transcribe_upload(request.files["audio"])
    intent_obj = detect_intent(stt["text"])
    return jsonify({"text": stt["text"], "intent": intent_obj, "meta": stt})

@app.post("/generate/word")
def generate_word():
    data = request.get_json(silent=True) or {}
    title = data.get("title", "Document")
    content = data.get("content", "")
    file_bytes = make_word_bytes(title, content)
    return send_file(
        BytesIO(file_bytes),
        mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        as_attachment=True,
        download_name="output.docx"
    )

@app.post("/generate/excel")
def generate_excel():
    data = request.get_json(silent=True) or {}
    sheet_name = data.get("sheet_name", "Sheet1")
    rows = data.get("rows", [["Task", "Status"], ["Example", "Open"]])
    file_bytes = make_excel_bytes(sheet_name, rows)
    return send_file(
        BytesIO(file_bytes),
        mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        as_attachment=True,
        download_name="output.xlsx"
    )

@app.post("/generate/ppt")
def generate_ppt():
    data = request.get_json(silent=True) or {}
    title = data.get("title", "Summary")
    bullets = data.get("bullets", ["Point 1", "Point 2"])
    file_bytes = make_ppt_bytes(title, bullets)
    return send_file(
        BytesIO(file_bytes),
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        as_attachment=True,
        download_name="output.pptx"
    )

# Local dev runner (keep for VS Code)
PORT = 5001

if __name__ == "__main__":
    PORT = int(os.getenv("PORT", "5001"))
    app.run(host="0.0.0.0", port=PORT, debug=True, use_reloader=False)

